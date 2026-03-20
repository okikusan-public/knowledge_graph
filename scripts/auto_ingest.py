#!/usr/bin/env python3
"""
GraphRAG自動インジェスト（プロジェクト共通版）

Usage:
  python auto_ingest.py upsert /path/to/file.md                    # デフォルトプロジェクト
  python auto_ingest.py upsert /path/to/file.md --project project_a  # プロジェクト指定
  python auto_ingest.py delete /path/to/file.md --project project_b
  GRAPHRAG_PROJECT=project_a python auto_ingest.py upsert /path/to/file.md
"""

import sys
import os
import uuid
import argparse
import requests
from neo4j import GraphDatabase

# 共通設定を読み込み
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from config import get_config


def get_driver(cfg):
    return GraphDatabase.driver(cfg.neo4j_uri, auth=cfg.neo4j_auth)


def gen_id():
    return str(uuid.uuid4())


def get_embeddings_batch(cfg, texts):
    """Embedding取得（バッチ対応）"""
    if not texts:
        return []
    prefixed = [f"passage: {t[:2000]}" if t.strip() else "passage: empty" for t in texts]
    resp = requests.post(cfg.embed_url, json={"inputs": prefixed}, timeout=60)
    resp.raise_for_status()
    return resp.json()["embeddings"]


def chunk_text(cfg, text):
    """テキストをチャンク分割"""
    char_size = int(cfg.chunk_size * cfg.char_per_token)
    char_overlap = int(cfg.chunk_overlap * cfg.char_per_token)

    if len(text) <= char_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = start + char_size
        chunk = text[start:end]

        if end < len(text):
            last_period = max(
                chunk.rfind("。"),
                chunk.rfind("\n\n"),
                chunk.rfind("\n"),
            )
            if last_period > char_size * 0.5:
                chunk = text[start:start + last_period + 1]
                end = start + last_period + 1

        chunks.append(chunk.strip())
        start = end - char_overlap

    return [c for c in chunks if c]


def extract_text(file_path):
    """ファイルからテキスト抽出"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext in (".md", ".txt", ".csv"):
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    elif ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(file_path)
            pages = [page.get_text() for page in doc]
            doc.close()
            return "\n\n".join(pages)
        except ImportError:
            print("  [warn] pymupdf not installed, skipping PDF", file=sys.stderr)
            return None

    elif ext in (".docx", ".doc"):
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path)
            parts = []

            # --- ヘッダー/フッター抽出 ---
            header_footer_texts = []
            for section in doc.sections:
                for hf in (section.header, section.footer):
                    if hf is not None:
                        hf_text = "\n".join(
                            p.text for p in hf.paragraphs if p.text.strip()
                        )
                        if hf_text.strip():
                            header_footer_texts.append(hf_text.strip())
            seen = set()
            for hf in header_footer_texts:
                if hf not in seen:
                    seen.add(hf)
                    parts.append(hf)

            # --- 段落抽出（見出し・リスト構造を保持） ---
            for para in doc.paragraphs:
                text = para.text
                if not text.strip():
                    continue
                style_name = (para.style.name or "") if para.style else ""

                if style_name.startswith("Heading"):
                    try:
                        level = int(style_name.split()[-1])
                    except (ValueError, IndexError):
                        level = 1
                    parts.append(f"{'#' * level} {text}")
                elif "List Bullet" in style_name:
                    indent_level = 0
                    try:
                        indent_level = int(style_name.split()[-1]) - 1
                    except (ValueError, IndexError):
                        pass
                    parts.append(f"{'  ' * indent_level}- {text}")
                elif "List Number" in style_name:
                    indent_level = 0
                    try:
                        indent_level = int(style_name.split()[-1]) - 1
                    except (ValueError, IndexError):
                        pass
                    parts.append(f"{'  ' * indent_level}1. {text}")
                else:
                    parts.append(text)

            # --- テーブル抽出（ネスト対応の再帰処理） ---
            def _extract_table(table):
                rows_text = []
                for row in table.rows:
                    cell_texts = []
                    for cell in row.cells:
                        nested_parts = []
                        cell_plain = cell.text.strip()
                        if cell.tables:
                            for p in cell.paragraphs:
                                if p.text.strip():
                                    nested_parts.append(p.text.strip())
                            for nested_table in cell.tables:
                                nested_parts.append(_extract_table(nested_table))
                            cell_texts.append(" ".join(nested_parts))
                        elif cell_plain:
                            cell_texts.append(cell_plain)
                    if cell_texts:
                        rows_text.append(" | ".join(cell_texts))
                return "\n".join(rows_text)

            for table in doc.tables:
                table_text = _extract_table(table)
                if table_text.strip():
                    parts.append(table_text)

            return "\n\n".join(parts)
        except ImportError:
            print("  [warn] python-docx not installed, skipping DOCX", file=sys.stderr)
            return None

    elif ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            from openpyxl.utils import range_boundaries
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheet_texts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text_parts = [f"=== シート: {sheet} ==="]

                # --- 結合セルの値を展開するルックアップを構築 ---
                merge_map = {}
                for merged_range in ws.merged_cells.ranges:
                    min_col, min_row, max_col, max_row = range_boundaries(str(merged_range))
                    origin_value = ws.cell(row=min_row, column=min_col).value
                    for r in range(min_row, max_row + 1):
                        for c in range(min_col, max_col + 1):
                            if (r, c) != (min_row, min_col):
                                merge_map[(r, c)] = origin_value

                # --- セル値取得ヘルパー（結合セル・数式Noneフォールバック対応） ---
                def cell_value(row_idx, col_idx):
                    if (row_idx, col_idx) in merge_map:
                        return merge_map[(row_idx, col_idx)]
                    val = ws.cell(row=row_idx, column=col_idx).value
                    if val is None:
                        return None
                    return val

                # --- 全行を読み取り ---
                if ws.max_row is None or ws.max_column is None:
                    continue
                rows_data = []
                for r in range(1, ws.max_row + 1):
                    row_vals = []
                    for c in range(1, ws.max_column + 1):
                        row_vals.append(cell_value(r, c))
                    rows_data.append(row_vals)

                # --- ヘッダー行の自動検出 ---
                header = None
                data_start = 0
                for i, row_vals in enumerate(rows_data):
                    non_none = [v for v in row_vals if v is not None]
                    if not non_none:
                        continue
                    if all(isinstance(v, str) for v in non_none):
                        header = row_vals
                        data_start = i + 1
                    else:
                        data_start = i
                    break

                # --- データ行の出力 ---
                for row_vals in rows_data[data_start:]:
                    if all(v is None for v in row_vals):
                        continue
                    if header is not None:
                        pairs = []
                        for col_name, val in zip(header, row_vals):
                            if val is None:
                                continue
                            label = str(col_name) if col_name is not None else ""
                            if label:
                                pairs.append(f"{label}: {val}")
                            else:
                                pairs.append(str(val))
                        if pairs:
                            text_parts.append(" | ".join(pairs))
                    else:
                        cells = [str(v) for v in row_vals if v is not None]
                        if cells:
                            text_parts.append(" | ".join(cells))

                sheet_texts.append("\n".join(text_parts))

            return "\n\n".join(sheet_texts)
        except ImportError:
            print("  [warn] openpyxl not installed, skipping Excel", file=sys.stderr)
            return None

    elif ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = [f"=== スライド {slide_num} ==="]
                # シェイプからテキスト抽出
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                slide_texts.append(para_text)
                    # テーブル内テキスト抽出
                    if shape.has_table:
                        for row in shape.table.rows:
                            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if cells:
                                slide_texts.append(" | ".join(cells))
                # スライドノート抽出
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_texts.append(f"[ノート] {notes_text}")
                text_parts.append("\n".join(slide_texts))
            return "\n\n".join(text_parts)
        except ImportError:
            print("  [warn] python-pptx not installed, skipping PPTX", file=sys.stderr)
            return None

    else:
        print(f"  [warn] Unsupported file type: {ext}", file=sys.stderr)
        return None


def cleanup_orphan_entities(driver):
    """どのChunkからもMENTIONSされていない孤立Entityを削除"""
    with driver.session() as session:
        result = session.run("""
            MATCH (e:Entity)
            WHERE NOT (e)<-[:MENTIONS]-(:Chunk)
            WITH e, e.name AS name
            DETACH DELETE e
            RETURN count(*) AS deleted, collect(name) AS names
        """)
        record = result.single()
        deleted = record["deleted"] if record else 0
        names = record["names"] if record else []
        if deleted > 0:
            print(f"  [cleanup] {deleted} orphaned entities removed: {', '.join(names[:10])}"
                  + ("..." if deleted > 10 else ""), file=sys.stderr)
        return deleted


def delete_document(driver, file_path):
    """ドキュメントと関連チャンクを削除"""
    with driver.session() as session:
        result = session.run("""
            MATCH (d:Document {source_path: $path})-[:HAS_CHUNK]->(c:Chunk)
            RETURN count(c) AS chunk_count
        """, path=file_path)
        record = result.single()
        chunk_count = record["chunk_count"] if record else 0

        # チャンクの関係を削除
        session.run("""
            MATCH (d:Document {source_path: $path})-[:HAS_CHUNK]->(c:Chunk)
            OPTIONAL MATCH (c)-[r]-()
            WHERE type(r) IN ['NEXT_CHUNK', 'MENTIONS', 'NEXT_STEP']
            DELETE r
        """, path=file_path)

        # チャンクを削除
        session.run("""
            MATCH (d:Document {source_path: $path})-[:HAS_CHUNK]->(c:Chunk)
            DETACH DELETE c
        """, path=file_path)

        # ドキュメントを削除
        session.run("""
            MATCH (d:Document {source_path: $path})
            DETACH DELETE d
        """, path=file_path)

        return chunk_count


def upsert_document(driver, cfg, file_path):
    """ドキュメントを追加または更新"""
    filename = os.path.basename(file_path)

    text = extract_text(file_path)
    if not text or not text.strip():
        print(f"  [skip] Empty or unreadable: {filename}", file=sys.stderr)
        return

    deleted_chunks = delete_document(driver, file_path)
    if deleted_chunks > 0:
        print(f"  [update] Removed old: {deleted_chunks} chunks", file=sys.stderr)

    chunks = chunk_text(cfg, text)

    # Embedding生成（ドキュメント全体 + 各チャンク）
    all_texts = [text[:3000]] + chunks
    embeddings = get_embeddings_batch(cfg, all_texts)
    doc_embedding = embeddings[0]
    chunk_embeddings = embeddings[1:]

    doc_id = gen_id()

    with driver.session() as session:
        # Documentノード作成
        session.run("""
            CREATE (d:Document {
                id: $id,
                title: $title,
                source_path: $path,
                file_type: $file_type,
                text_length: $text_length,
                chunk_count: $chunk_count,
                embedding: $embedding,
                auto_ingested: true,
                created_at: datetime()
            })
        """, id=doc_id, title=filename, path=file_path,
             file_type=os.path.splitext(file_path)[1].lstrip("."),
             text_length=len(text),
             chunk_count=len(chunks), embedding=doc_embedding)

        # Chunkノード作成 + HAS_CHUNK関係
        chunk_ids = []
        for i, (chunk, emb) in enumerate(zip(chunks, chunk_embeddings)):
            chunk_id = gen_id()
            chunk_ids.append(chunk_id)

            session.run("""
                MATCH (d:Document {id: $doc_id})
                CREATE (c:Chunk {
                    id: $chunk_id,
                    text: $text,
                    chunk_index: $index,
                    token_estimate: $tokens,
                    embedding: $embedding
                })
                CREATE (d)-[:HAS_CHUNK]->(c)
            """, doc_id=doc_id, chunk_id=chunk_id, text=chunk,
                 index=i, tokens=int(len(chunk) / cfg.char_per_token),
                 embedding=emb)

        # NEXT_CHUNK関係作成
        for i in range(len(chunk_ids) - 1):
            session.run("""
                MATCH (c1:Chunk {id: $id1}), (c2:Chunk {id: $id2})
                CREATE (c1)-[:NEXT_CHUNK]->(c2)
            """, id1=chunk_ids[i], id2=chunk_ids[i + 1])

        # 既存エンティティとのMENTIONS関係作成
        session.run("""
            MATCH (d:Document {id: $doc_id})-[:HAS_CHUNK]->(c:Chunk)
            MATCH (e:Entity)
            WHERE c.text CONTAINS e.name
            MERGE (c)-[:MENTIONS]->(e)
        """, doc_id=doc_id)

    print(f"  [done] {filename}: {len(chunks)} chunks ({cfg.project})", file=sys.stderr)


def main():
    parser = argparse.ArgumentParser(description="GraphRAG Auto Ingest")
    parser.add_argument("action", choices=["upsert", "delete"])
    parser.add_argument("file_path")
    parser.add_argument("--project", "-p", default=None,
                        help="プロジェクト名 (config.py の PROJECTS キー)")
    args = parser.parse_args()

    cfg = get_config(args.project)
    print(f"  [config] {cfg}", file=sys.stderr)

    driver = get_driver(cfg)
    try:
        if args.action == "upsert":
            upsert_document(driver, cfg, args.file_path)
            cleanup_orphan_entities(driver)
        elif args.action == "delete":
            deleted = delete_document(driver, args.file_path)
            print(f"  [deleted] {os.path.basename(args.file_path)}: {deleted} chunks removed",
                  file=sys.stderr)
            cleanup_orphan_entities(driver)
    finally:
        driver.close()


if __name__ == "__main__":
    main()
